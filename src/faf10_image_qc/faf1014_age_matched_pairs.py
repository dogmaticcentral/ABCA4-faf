#!/usr/bin/env python
"""
Finds age-matched pairs of images from control and non-control groups.
Both images must be marked as 'usable'.
"""

import sys
import os
import tempfile
from typing import List, Tuple, Set
from collections import defaultdict
from PIL import Image

# Peewee Models & Util
from models.abca4_faf_models import FafImage, Case
from utils.db_utils import db_connect
from faf00_settings import global_db_proxy

# Reports
try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from utils.reports import pptx_to_pdf
except ImportError:
    print("WARNING: python-pptx not installed. PPTX generation will fail.", file=sys.stderr)

# ---------------------------------------------------------------------------
# Database & matching Logic
# ---------------------------------------------------------------------------

def get_usable_images() -> List[FafImage]:
    """Fetches usable images with their case info."""
    # Ensure DB connection
    if global_db_proxy.obj is None:
         db = db_connect()
    else:
         db = global_db_proxy
         db.connect(reuse_if_open=True)

    # Join with Case to filter/access control status eagerly
    query = (FafImage
             .select(FafImage, Case)
             .join(Case)
             .where(FafImage.usable == True))
    
    return list(query)

def match_images(patients: List[FafImage], controls: List[FafImage]) -> List[Tuple[FafImage, FafImage, float]]:
    """
    Greedy matching of patients to controls based on age.
    Returns list of (patient_img, control_img, age_diff).
    Strategy: Calculate all pair differences, sort by smallest difference, pick distinct pairs.
    """
    potential_pairs = []
    for p in patients:
        for c in controls:
            if p.age_acquired is not None and c.age_acquired is not None:
                diff = abs(p.age_acquired - c.age_acquired)
                potential_pairs.append((diff, p, c))
    
    # Sort by smallest difference first
    potential_pairs.sort(key=lambda x: x[0])
    
    matched_pairs = []
    matched_p_ids: Set[int] = set()
    matched_c_ids: Set[int] = set()
    
    for diff, p, c in potential_pairs:
        if p.id not in matched_p_ids and c.id not in matched_c_ids:
            matched_pairs.append((p, c, diff))
            matched_p_ids.add(p.id)
            matched_c_ids.add(c.id)
            
    return matched_pairs

def get_visit_images(case_id: int, age: float) -> dict[str, str]:
    """
    Returns a dict {'OD': path, 'OS': path} for a given case and age.
    """
    # Use a small epsilon for float comparison just in case
    # Note: peewee SQL generation might be safer with python-side filter if not too many
    # But FafImage.select().where(...) is standard
    
    siblings = (FafImage
                .select()
                .where(
                    (FafImage.case_id == case_id) & 
                    (FafImage.usable == True) &
                    (FafImage.age_acquired >= age - 0.01) & 
                    (FafImage.age_acquired <= age + 0.01)
                ))
    
    res = {}
    for img in siblings:
        res[img.eye] = img.image_path
    return res

# ---------------------------------------------------------------------------
# PPTX / Report Logic
# ---------------------------------------------------------------------------

def ensure_compatible_image(image_path: str) -> str:
    """
    Ensures image is compatible with PPTX (JPG, PNG). 
    If TIFF, converts to temporary PNG.
    Returns path to usable image.
    """
    if not image_path or not os.path.exists(image_path):
        return None
        
    lower = image_path.lower()
    if lower.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.gif')):
        return image_path
    
    # Needs conversion (likely TIFF)
    try:
        with Image.open(image_path) as img:
            # Create a temp file
            fd, temp_path = tempfile.mkstemp(suffix='.png')
            os.close(fd)
            # Convert to RGB if necessary (e.g. 16-bit TIFF or CMYK)
            if img.mode not in ('RGB', 'RGBA'):
                img = img.convert('RGB')

            # Crop to 1/3 width × 1/2 height around center
            width, height = img.size
            new_width = width // 3
            new_height = height // 2

            left = (width - new_width) // 2
            top = (height - new_height) // 2
            right = left + new_width
            bottom = top + new_height

            img = img.crop((left, top, right, bottom))

            img.save(temp_path, format='PNG')
        return temp_path
    except Exception as e:
        print(f"WARN: Could not convert {image_path}: {e}", file=sys.stderr)
        return None

def add_slide(prs: Presentation, 
              patient_alias: str, patient_age: float, patient_imgs: dict, 
              control_alias: str, control_age: float, control_imgs: dict):
    
    title_slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = f"P: {patient_alias} ({patient_age:.1f}y) vs C: {control_alias} ({control_age:.1f}y)"
    
    # Layout constants
    margin = Inches(0.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Available area
    content_width = slide_width - 2 * margin
    content_height = slide_height - Inches(1.5) # Reserve space for title
    
    # 2 rows (Patient, Control), 2 columns (OD, OS)
    # Row heights
    row_h = content_height / 2
    col_w = content_width / 2
    
    y_start = Inches(1.2)
    
    # Helper to place image
    def place_image(path, r, c, label):
        if not path: return
        
        real_path = ensure_compatible_image(path)
        if not real_path: return
        
        x = margin + c * col_w
        y = y_start + r * row_h
        
        # Calculate fitting (naive centered fit)
        # We define a box for the image
        box_w = col_w * 0.95
        box_h = row_h * 0.8
        
        # Add picture
        try:
            pic = slide.shapes.add_picture(real_path, x, y, width=box_w)
            # Adjust height if it exceeds box_h (maintain aspect ratio)
            if pic.height > box_h:
                pic.height = box_h
            
            # Center horizontally in the slot
            # Note: pic.left/top are updated after resize if we set them, but let's re-center
            pic.left = int(x + (box_w - pic.width)/2)
                
            # Add label
            tx = slide.shapes.add_textbox(x, y + box_h + Pt(5), box_w, Pt(20))
            tf = tx.text_frame
            tf.text = label
            tf.paragraphs[0].font.size = Pt(12)
            
            # If temp file, we could clean up, but we might reuse it if deduping failed.
            # For simplicity, we leave temp files to OS cleanup or delete later.
            if real_path != path:
                try:
                    os.remove(real_path)
                except:
                    pass
        except Exception as e:
            print(f"ERROR placing image {path}: {e}", file=sys.stderr)

    place_image(patient_imgs.get('OD'), 0, 0, f"Patient OD")
    place_image(patient_imgs.get('OS'), 0, 1, f"Patient OS")
    place_image(control_imgs.get('OD'), 1, 0, f"Control OD")
    place_image(control_imgs.get('OS'), 1, 1, f"Control OS")


def main():
    images = get_usable_images()
    
    patients = []
    controls = []
    
    for img in images:
        if img.case_id.is_control:
            controls.append(img)
        else:
            patients.append(img)
            
    # Print stats to stderr so stdout can be piped clean
    print(f"INFO: Found {len(patients)} usable patient images and {len(controls)} usable control images.", file=sys.stderr)
    
    pairs = match_images(patients, controls)
    
    print(f"INFO: Matched {len(pairs)} pairs.", file=sys.stderr)

    # Sort pairs by patient age (ascending)
    pairs.sort(key=lambda x: x[0].age_acquired)
    
    # TSV Output
    print("Patient_Alias\tPatient_Age\tControl_Alias\tControl_Age\tAge_Diff")
    for p, c, diff in pairs:
        print(f"{p.case_id.alias}\t{p.age_acquired}\t{c.case_id.alias}\t{c.age_acquired}\t{diff:.4f}")

    # --- PPTX Generation ---
    out_name = "age_matched_pairs.pptx"
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Age-Matched Pairs"
    slide.placeholders[1].text = f"Total Pairs: {len(pairs)}"
    
    # Deduplicate by (PatientCase, ControlCase, PatientAge) to avoid redundant slides
    seen_visits = set()
    
    for p, c, diff in pairs:
        p_case = p.case_id
        c_case = c.case_id
        
        # Key for the visit/match
        key = (p_case.id, p.age_acquired, c_case.id)
        
        if key in seen_visits:
            continue
        seen_visits.add(key)
        
        # Get all images for these visits
        p_imgs = get_visit_images(p_case.id, p.age_acquired)
        c_imgs = get_visit_images(c_case.id, c.age_acquired)
        
        add_slide(prs, 
                  p_case.alias, p.age_acquired, p_imgs,
                  c_case.alias, c.age_acquired, c_imgs)
                  
    prs.save(out_name)
    print(f"INFO: Saved PPTX to {out_name}", file=sys.stderr)
    
    try:
        pptx_to_pdf(out_name)
        print(f"INFO: Converted to PDF.", file=sys.stderr)
    except Exception as e:
        print(f"WARN: PDF conversion failed: {e}", file=sys.stderr)

if __name__ == "__main__":
    main()
