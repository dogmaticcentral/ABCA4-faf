_copyright__ = """

    Copyright 2024 Ivana Mihalek

    Licensed under Creative Commons Attribution-NonCommercial 4.0 International Public License:
    You may obtain a copy of the License at https://creativecommons.org/licenses/by-nc/4.0/
    
    The License is noncommercial - you may not use this material for commercial purposes.

"""
__license__ = "CC BY-NC 4.0"

import os
from datetime import datetime
from pathlib import Path

from faf00_settings import SOFFICE, WORK_DIR
from pptx import Presentation
from pptx.util import Pt

from utils.conventions import construct_report_filepath
from utils.processes import run_subprocess
from utils.utils import is_nonempty_file, is_runnable, scream, shrug


def pptx_to_pdf(pptx_filepath: Path | str, keep_pptx: bool = False) -> Path:
    if not SOFFICE or not is_runnable(SOFFICE):
        shrug("I cannot convert to pdf unless soffice is installed and runnable. Keeping the pttx.")
        return pptx_filepath
    pptx_filepath = Path(pptx_filepath)
    pdf_filepath =  pptx_filepath.parent / (pptx_filepath.stem + ".pdf")
    outdir = pdf_filepath.parent  # soffice will take the stem and attach the pdf suffix
    run_subprocess(f"{SOFFICE} soffice --headless --convert-to pdf {pptx_filepath} --outdir {outdir}")
    if is_nonempty_file(pdf_filepath):
        print(f"wrote {pdf_filepath}")
        if not keep_pptx: os.unlink(pptx_filepath)
        return pdf_filepath
    else:
        scream("creating {pdf_filepath} failed")
        exit()


######################################################################
def make_paired_slides(img_filepaths: dict, name_stem: str, title: str = "") -> Path:
    pptx_filepath = construct_report_filepath(WORK_DIR, name_stem,  "pptx")
    prs = Presentation()

    # the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slides_title = slide.shapes.title
    subtitle = slide.placeholders[1]

    slides_title.text = title if title else name_stem.capitalize()
    subtitle.text = datetime.today().strftime('%B %d, %Y')
    title_only_slide_layout = prs.slide_layouts[5]

    # the payload
    left = prs.slide_width * 0.05
    top = prs.slide_height * 0.20
    img_height = prs.slide_height * 0.35
    img_width  = prs.slide_width  * 0.6
    for alias, image_pairs in img_filepaths.items():
        for [right_eye_img, left_eye_img] in image_pairs:
            slide = prs.slides.add_slide(title_only_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = alias
            title_shape.text_frame.paragraphs[0].font.size = Pt(28)
            if right_eye_img:  # in some cases we can have only one eye
                slide.shapes.add_picture(right_eye_img, left, top, height=img_height)
                text_frame_top = slide.shapes.add_textbox(left + img_width, top, width=img_width//2, height=img_height//4)
                text_frame_top.text = Path(right_eye_img).stem
            if left_eye_img:
                slide.shapes.add_picture(left_eye_img,  left, top + img_height, height=img_height)
                text_frame_bottom = slide.shapes.add_textbox(left + img_width, top + img_height, width=img_width//2, height=img_height//4)
                text_frame_bottom.text = Path(left_eye_img).stem

    prs.save(pptx_filepath)
    return pptx_filepath


def make_paired_pdf(img_filepaths: dir, name_stem: str, title: str = "", keep_pptx=False) -> Path | None:
    if not SOFFICE or not is_runnable(SOFFICE):
        shrug(f"I dont know how to create pdf without soffice installed and runnable. You can try creating pptx instead.")
        return None

    slides_filepath = make_paired_slides(img_filepaths, name_stem, title)
    print(f"Created {slides_filepath}. Converting to pdf.")
    pdf_filepath = pptx_to_pdf(slides_filepath, keep_pptx=keep_pptx)
    return pdf_filepath
