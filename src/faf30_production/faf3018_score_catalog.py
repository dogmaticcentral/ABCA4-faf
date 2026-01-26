#!/usr/bin/env python

"""
    © 2024-2026 Ivana Mihalek ivana.mihalek@gmail.com

    Licensed under Creative Commons Attribution-NonCommercial 4.0 International Public License:
    You may obtain a copy of the License at https://creativecommons.org/licenses/by-nc/4.0/

    The License is noncommercial - you may not use this material for commercial purposes.

"""

from statistics import mean

import peewee
from pptx import Presentation
from pptx.util import Pt

from faf00_settings import WORK_DIR, USE_AUTO, DATABASES, DEBUG
from models.abca4_faf_models import FafImage
from models.abca4_results import Score
from utils.conventions import construct_workfile_path
from utils.db_utils import db_connect
from faf00_settings import global_db_proxy
from utils.reports import pptx_to_pdf
from utils.utils import is_nonempty_file, shrug


def rows_w_avg_score_from_db() -> list[list]:
    rows = []  # [alias, age, avg_score, orig_images: dict]
    we_are_using_postgres = DATABASES["default"]["ENGINE"] == 'peewee.postgres'
    query = FafImage.select(
        FafImage.case_id,
        FafImage.age_acquired,
        peewee.fn.array_agg(FafImage.id).alias("img_ids") if we_are_using_postgres
        else peewee.fn.GROUP_CONCAT(FafImage.id).alias("img_ids"),
    ).group_by(FafImage.case_id, FafImage.age_acquired)

    for faf_img in query:
        alias = faf_img.case_id.alias

        if isinstance(faf_img.img_ids, int):
            pair_imgs = [faf_img.img_ids]  # this one is actually missing its pair
        else:
            pairs = faf_img.img_ids if we_are_using_postgres else faf_img.img_ids.split(",")
            pair_imgs = [int(i) for i in pairs]

        try:
            if USE_AUTO:
                avg_score = mean([Score.select().where(Score.faf_image_id==i).get().pixel_score_auto for i in pair_imgs])
            else:
                avg_score = mean([Score.select().where(Score.faf_image_id==i).get().pixel_score for i in pair_imgs])
        except Exception:
            shrug(f"one or more images from the pair {pair_imgs} has no score present in the db")
            continue
        age_images_acquired = faf_img.age_acquired
        image_paths = {(entry := FafImage.get_by_id(i)).eye: entry.image_path for i in pair_imgs}
        rows.append([alias, age_images_acquired, avg_score, image_paths])
    rows_sorted = sorted(rows, key=lambda row: row[2])
    return rows_sorted


def add_score_progression_slide(prs: Presentation, row: list):
    [alias, age, avg_score, orig_images] = row
    if not orig_images:
        return
    if DEBUG: print(alias, age, f"{avg_score:.2f}")

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = f"{alias.replace('_', ' ')}, age {age}, avg score: {avg_score:.2f}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(22)

    left = prs.slide_width * 0.05
    top = prs.slide_height * 0.3
    img_width = (prs.slide_width - 2 * left - left / 2) / 2

    for eye, orig_image in orig_images.items():
        comp_path = str(construct_workfile_path(WORK_DIR, orig_image, alias, "composite", "png"))
        composite_path = comp_path if is_nonempty_file(comp_path) else None
        score_png = str(construct_workfile_path(WORK_DIR, orig_image, alias, "pixel_score", "png"))
        pixel_score_path = score_png if is_nonempty_file(score_png) else None
        if eye == "OD":
            if composite_path:
                slide.shapes.add_picture(composite_path, left, top, width=img_width)
            if pixel_score_path:
                slide.shapes.add_picture(pixel_score_path, left, top + prs.slide_height / 3, width=img_width)
        if eye == "OS":
            if composite_path:
                slide.shapes.add_picture(composite_path, left + img_width + left / 2, top, width=img_width)
            if pixel_score_path:
                x_coord = left + img_width + left / 2
                y_coord = top  + prs.slide_height / 3
                slide.shapes.add_picture(pixel_score_path, x_coord, y_coord, width=img_width)

    return


def main():

    out_name = "score_progression.pptx"

    # get the score info from the database
    if global_db_proxy.obj is None:
         db = db_connect()
    else:
         db = global_db_proxy
         db.connect(reuse_if_open=True)
    rows = rows_w_avg_score_from_db()
    
    if not db.is_closed():
        db.close()

    # initialize the pttx presentation
    prs = Presentation()

    # the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Abca4 FAF cases - avg pixel score"
    subtitle.text = "Ivana, March 2024"

    for row in rows:
        add_score_progression_slide(prs, row)


    prs.save(out_name)
    print(f"Saved as {out_name}. Converting to pdf. (The pptx will be removed.)")
    pptx_to_pdf(out_name)


########################################
if __name__ == "__main__":
    main()
