#!/usr/bin/env python

"""
    © 2024 Ivana Mihalek ivana.mihalek@gmail.com

    Licensed under Creative Commons Attribution-NonCommercial 4.0 International Public License:
    You may obtain a copy of the License at https://creativecommons.org/licenses/by-nc/4.0/

    The License is noncommercial - you may not use this material for commercial purposes.

"""


import os

from pptx import Presentation
from pptx.util import Pt

from faf00_settings import WORK_DIR, USE_AUTO
from models.abca4_results import Score
from utils.conventions import construct_workfile_path
from utils.db_utils import db_connect
from utils.reports import pptx_to_pdf
from utils.utils import is_nonempty_file, shrug


def rows_from_db() -> list[list]:
    rows = []

    for score in Score.select():
        faf_img = score.faf_image_id
        case = faf_img.case_id
        rows.append([case.alias, faf_img.age_acquired, faf_img.eye,
                     score.pixel_score_auto if USE_AUTO else score.pixel_score,
                     faf_img.image_path])

    return rows


def rows_ordered_by_alias_and_age(rows):

    controls = []
    tmp_dict = {}
    for row in rows:
        [alias, age, eye] = row[:3]
        if "control" in alias.lower():
            controls.append(row)
            continue
        age = float(age)
        if alias not in tmp_dict: tmp_dict[alias] = {}
        if age not in tmp_dict[alias]: tmp_dict[alias][age] = {}
        tmp_dict[alias][age][eye] = row

    sorted_rows = []
    for alias in sorted(tmp_dict.keys()):
        for age in sorted(tmp_dict[alias].keys()):
            for eye in sorted(tmp_dict[alias][age].keys()):
                sorted_rows.append(tmp_dict[alias][age][eye])

    return sorted_rows + controls


def add_catalog_slide(prs: Presentation, row: list):

    [alias, age, eye, score, original_img_path] = row

    composite_png = construct_workfile_path(WORK_DIR, original_img_path, alias, "composite", 'png')
    if USE_AUTO:
        bg_hist_png   = construct_workfile_path(WORK_DIR, original_img_path, alias, "auto_bg_histogram" , 'png')
        if not is_nonempty_file(bg_hist_png):
            shrug(f"{bg_hist_png} does not exist (or may be empty) - falling back on the manual selection.")
            bg_hist_png   = construct_workfile_path(WORK_DIR, original_img_path, alias, "bg_histogram", 'png')
    else:
        bg_hist_png   = construct_workfile_path(WORK_DIR, original_img_path, alias, "bg_histogram", 'png')
    hist_png      = construct_workfile_path(WORK_DIR, original_img_path, alias, "roi_histogram", 'png')
    score_png     = construct_workfile_path(WORK_DIR, original_img_path, alias, "pixel_score", 'png')

    for img in [composite_png, bg_hist_png, hist_png, score_png]:
        if not os.path.exists(img):
            print(f"\t{img} not found")
            return
        if os.path.getsize(img) == 0:
            print(f"\t{img} seems to be empty")
            return

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = f"{alias.replace('_', ' ')}, {eye}, {age}, score: {int(score)}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(16)

    left = prs.slide_width*0.05
    top  = prs.slide_height*0.20
    img_height = prs.slide_height*0.4

    slide.shapes.add_picture(str(composite_png),  left,                top, height=img_height)
    slide.shapes.add_picture(str(score_png),     left,                     top + img_height, height=img_height)
    slide.shapes.add_picture(str(bg_hist_png),   left + prs.slide_width/2, top, height=img_height)
    slide.shapes.add_picture(str(hist_png), left + prs.slide_width/2, top + img_height, height=img_height)

    return


def main():

    # TODO make this an argv
    out_name = 'catalog.pptx'

    # get the score info from the database
    db = db_connect()
    rows = rows_from_db()
    db.close()

    # initialize the pttx presentation
    prs = Presentation()

    # the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Abca4 FAF image catalog"
    subtitle.text = "Ivana, Sept 2024"

    for row in rows_ordered_by_alias_and_age(rows):
        add_catalog_slide(prs, row)

    prs.save(out_name)
    print(f"Saved as {out_name}. Converting to pdf. (The pptx will be removed.)")
    pptx_to_pdf(out_name)


########################################
if __name__ == "__main__":
    main()
